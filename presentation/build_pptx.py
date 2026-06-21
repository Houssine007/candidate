# -*- coding: utf-8 -*-
"""
Génère la présentation de soutenance RecruitPRO (.pptx).
Style : académique, minimaliste, sobre. Blanc / bleu marine / bleu ONCF / gris.
Texte réduit à l'écran ; le détail oral est placé dans les NOTES de chaque slide.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---- Palette ---------------------------------------------------------------
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
INK        = RGBColor(0x0F, 0x2A, 0x4A)   # bleu marine (titres, texte fort)
BLUE       = RGBColor(0x1C, 0x5F, 0xA8)   # bleu ONCF (accent)
BLUE_DK    = RGBColor(0x14, 0x3A, 0x66)
GRAY       = RGBColor(0x5B, 0x66, 0x71)   # texte secondaire
GRAY_SOFT  = RGBColor(0x8A, 0x93, 0x9E)
LINE       = RGBColor(0xD9, 0xDE, 0xE5)   # filets
PANEL      = RGBColor(0xF5, 0xF7, 0xF9)   # fonds clairs
BLUE_LT    = RGBColor(0xE9, 0xF1, 0xFA)   # bleu très clair

F_TITLE = "Segoe UI Semibold"
F_BODY  = "Segoe UI"
F_LIGHT = "Segoe UI Light"

EMU_IN = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

# ---- Helpers ---------------------------------------------------------------
def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    return s

def _set_run(r, text, font=F_BODY, size=18, color=INK, bold=False, italic=False):
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color

def textbox(s, l, t, w, h, lines, font=F_BODY, size=18, color=INK, bold=False,
            align=PP_ALIGN.LEFT, italic=False, space_after=6, line_spacing=1.05,
            anchor=MSO_ANCHOR.TOP):
    """lines: str OR list[str] OR list[list[(text, overrides)]]"""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        if isinstance(ln, str):
            _set_run(p.add_run(), ln, font, size, color, bold, italic)
        else:  # list of (text, dict)
            for seg in ln:
                txt, ov = seg
                _set_run(p.add_run(), txt, ov.get("font", font), ov.get("size", size),
                         ov.get("color", color), ov.get("bold", bold), ov.get("italic", italic))
    return tb

def rect(s, l, t, w, h, fill=None, line_color=None, line_w=1.0, rounded=False, radius=0.08):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color; shp.line.width = Pt(line_w)
    if rounded:
        try: shp.adjustments[0] = radius
        except Exception: pass
    return shp

def shape_text(shp, lines, font=F_BODY, size=15, color=INK, bold=False, align=PP_ALIGN.CENTER,
               space_after=2, line_spacing=1.0):
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in (tf.margin_left, ):
        pass
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04);  tf.margin_bottom = Inches(0.04)
    if isinstance(lines, str): lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        if isinstance(ln, str):
            _set_run(p.add_run(), ln, font, size, color, bold)
        else:
            for seg in ln:
                txt, ov = seg
                _set_run(p.add_run(), txt, ov.get("font", font), ov.get("size", size),
                         ov.get("color", color), ov.get("bold", bold))

def arrow(s, x1, y1, x2, y2, color=BLUE, width=1.75):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    ln = cn.line._get_or_add_ln()
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    cn.shadow.inherit = False
    return cn

def hline(s, l, t, w, color=LINE, width=1.0):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(l), Inches(t), Inches(l + w), Inches(t))
    cn.line.color.rgb = color; cn.line.width = Pt(width); cn.shadow.inherit = False
    return cn

def header(s, title, kicker=None):
    if kicker:
        textbox(s, 0.92, 0.55, 11.5, 0.35, kicker.upper(), font=F_BODY, size=11.5,
                color=BLUE, bold=True)
        ty = 0.86
    else:
        ty = 0.7
    textbox(s, 0.9, ty, 11.5, 0.9, title, font=F_TITLE, size=30, color=INK, bold=False)
    rect(s, 0.92, ty + 0.78, 0.62, 0.045, fill=BLUE)   # accent rule
    return ty + 1.15

def footer(s, idx, label):
    hline(s, 0.9, 7.0, SW - 1.8, color=LINE, width=0.75)
    textbox(s, 0.9, 7.06, 8.0, 0.3, "RecruitPRO — " + label, font=F_BODY, size=9.5, color=GRAY_SOFT)
    textbox(s, SW - 1.7, 7.06, 0.8, 0.3, f"{idx:02d} / 15", font=F_BODY, size=9.5,
            color=GRAY_SOFT, align=PP_ALIGN.RIGHT)

def bullets(s, items, top, left=0.95, width=11.4, size=19, gap=0.86, color=INK,
            marker=BLUE, sub_color=GRAY):
    """items: list of str OR (main, sub)"""
    y = top
    for it in items:
        rect(s, left, y + 0.06, 0.12, 0.12, fill=marker)   # square marker
        if isinstance(it, tuple):
            main, sub = it
            textbox(s, left + 0.32, y - 0.08, width - 0.32, 0.5, main, font=F_BODY, size=size,
                    color=color, bold=False)
            textbox(s, left + 0.32, y + 0.30, width - 0.32, 0.4, sub, font=F_BODY, size=13.5,
                    color=sub_color)
        else:
            textbox(s, left + 0.32, y - 0.08, width - 0.32, 0.5, it, font=F_BODY, size=size, color=color)
        y += gap

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text.strip()

# ===========================================================================
# SLIDE 1 — Titre
# ===========================================================================
s = slide()
rect(s, 0, 0, SW, 0.22, fill=INK)                      # bandeau supérieur fin
rect(s, 0, 0, 0.22, SH, fill=BLUE)                     # filet latéral accent
textbox(s, 1.1, 2.15, 11.0, 1.3, "RecruitPRO", font=F_LIGHT, size=66, color=INK)
textbox(s, 1.12, 3.45, 10.8, 0.7,
        "Vers un référentiel de compétences vivant", font=F_TITLE, size=24, color=BLUE)
hline(s, 1.13, 4.25, 4.6, color=LINE, width=1.0)
textbox(s, 1.12, 4.45, 11.0, 0.5,
        "Plateforme RH intégrée  ·  ATS · GPEC · LMS", font=F_BODY, size=15, color=GRAY)
textbox(s, 1.12, 5.7, 11.0, 1.2, [
    "Mémoire de Master — Génie Logiciel",
    "Cadre : ONCF — Institut de Formation et de Développement des Compétences (IFDC)",
], font=F_BODY, size=13.5, color=GRAY, space_after=4)
textbox(s, 1.12, 6.65, 11.0, 0.5, [[
    ("Présenté par : Houssine Mir", {"color": INK, "bold": True}),
    ("     ·     Sous la direction de : Pr. [à compléter]     ·     2026", {"color": GRAY}),
]], font=F_BODY, size=12.5)
notes(s, """
Bonjour. Monsieur le Président, Mesdames et Messieurs les membres du jury, je vous remercie de votre présence.
Je vais vous présenter RecruitPRO, une plateforme RH intégrée conçue dans le cadre de l'IFDC de l'ONCF.
Le fil conducteur de tout mon travail tient en une idée : faire vivre le référentiel de compétences.
La présentation dure une quinzaine de minutes ; je terminerai par une démonstration et serai à votre disposition pour vos questions.
""")

# ===========================================================================
# SLIDE 2 — Contexte
# ===========================================================================
s = slide()
y = header(s, "Contexte", "01 · Cadre")
bullets(s, [
    "L'ONCF investit fortement dans la formation.",
    "Les métiers et les compétences évoluent en continu.",
    "Recrutement, mobilité et formation restent séparés.",
    "La donnée « compétence » demeure cloisonnée.",
], top=y + 0.15, size=20, gap=0.95)
footer(s, 2, "Contexte")
notes(s, """
L'ONCF est une grande organisation qui investit massivement dans le développement de ses collaborateurs.
Mais dans la plupart des organisations, trois mondes coexistent sans se parler : le recrutement, la gestion de carrière, et la formation.
Chacun a son outil, sa base, sa logique. La compétence — qui est pourtant l'actif central des RH — circule mal entre eux.
Mon point de départ est donc simple : l'enjeu n'est pas l'outil, c'est la compétence, et sa capacité à rester à jour.
""")

# ===========================================================================
# SLIDE 3 — Problématique
# ===========================================================================
s = slide()
header(s, "Problématique", "02 · Le problème")
textbox(s, 0.9, 2.5, 11.5, 1.4,
        "Le référentiel de compétences est figé.", font=F_LIGHT, size=40, color=INK)
bullets(s, [
    "Saisi une fois, rarement mis à jour.",
    "Déconnecté de ce qui se passe sur le terrain.",
    "Il se périme plus vite qu'on ne l'entretient.",
], top=4.05, size=18, gap=0.62)
textbox(s, 0.9, 6.1, 11.5, 0.6, [[
    ("Question de recherche : ", {"color": GRAY}),
    ("comment faire vivre ce référentiel ?", {"color": BLUE, "bold": True}),
]], font=F_BODY, size=18)
footer(s, 3, "Problématique")
notes(s, """
Voici le cœur du problème. Dans presque tous les systèmes RH, le référentiel de compétences est un document figé.
On le saisit une fois, à l'embauche ou lors d'un audit, puis on l'oublie. Il ne reflète plus ni les formations suivies, ni les mobilités, ni l'expérience acquise.
Résultat : il se périme plus vite qu'on ne parvient à l'entretenir manuellement. On pilote les compétences avec une photo, alors qu'il faudrait un film.
Ma question de recherche est donc : comment faire pour que ce référentiel se mette à jour tout seul, à partir de la vie réelle de l'organisation ?
""")

# ===========================================================================
# SLIDE 4 — Objectifs
# ===========================================================================
s = slide()
y = header(s, "Objectifs", "03 · Cap")
bullets(s, [
    ("Unifier ATS, RH/GPEC et LMS.", "Une seule identité, une seule vérité."),
    ("Mettre à jour le référentiel automatiquement.", "À partir d'événements réels."),
    ("Évaluer l'adéquation et le potentiel.", "Recruter le présent, anticiper l'avenir."),
    ("Assurer la portabilité candidat → employé.", "Le profil suit la personne."),
], top=y + 0.1, size=19, gap=1.02)
footer(s, 4, "Objectifs")
notes(s, """
De cette question découlent quatre objectifs.
Premièrement, unifier les trois mondes — recrutement, RH, formation — sous une seule identité numérique.
Deuxièmement, et c'est le plus important, mettre à jour le référentiel automatiquement, à partir d'événements réels et non d'une saisie manuelle.
Troisièmement, ne pas se limiter à « est-ce que le candidat correspond ? » mais aussi « peut-il évoluer ? ».
Quatrièmement, garantir que le profil de compétences suive la personne, du statut de candidat à celui d'employé.
""")

# ===========================================================================
# SLIDE 5 — État de l'art
# ===========================================================================
s = slide()
y = header(s, "État de l'art", "04 · L'existant")
cols = [
    ("ATS", "Recrute, puis oublie le profil."),
    ("LMS", "Forme, sans relier au poste."),
    ("SIRH", "Stocke, ne fait pas évoluer."),
]
cw, gap_c, x0, top = 3.5, 0.45, 0.95, y + 0.15
for i, (t, d) in enumerate(cols):
    x = x0 + i * (cw + gap_c)
    card = rect(s, x, top, cw, 2.0, fill=PANEL, line_color=LINE, line_w=1.0, rounded=True, radius=0.06)
    textbox(s, x + 0.25, top + 0.28, cw - 0.5, 0.6, t, font=F_TITLE, size=22, color=INK)
    hline(s, x + 0.27, top + 1.0, 0.5, color=BLUE, width=2.0)
    textbox(s, x + 0.25, top + 1.15, cw - 0.5, 0.8, d, font=F_BODY, size=15, color=GRAY)
textbox(s, 0.9, top + 2.5, 11.5, 0.7, [[
    ("Le constat : ", {"color": GRAY}),
    ("aucun ne fait vivre le référentiel.", {"color": BLUE, "bold": True}),
]], font=F_BODY, size=20)
footer(s, 5, "État de l'art")
notes(s, """
Qu'offre l'existant ? On peut le résumer en trois familles d'outils.
Les ATS gèrent le recrutement, mais une fois la personne embauchée, le profil de compétences est abandonné.
Les LMS dispensent de la formation, mais sans la relier au poste ni au besoin réel de l'organisation.
Les SIRH stockent les données du personnel, mais comme une archive : ils ne font pas évoluer la compétence.
Le constat est net : chacun fait bien sa part, mais aucun ne fait vivre le référentiel en reliant ces trois moments. C'est précisément l'espace que RecruitPRO occupe.
""")

# ===========================================================================
# SLIDE 6 — Architecture générale (diagramme)
# ===========================================================================
s = slide()
header(s, "Architecture générale", "05 · Vue système")
# bandeau SSO
sso = rect(s, 1.6, 1.95, 10.1, 0.7, fill=INK, rounded=True, radius=0.5)
shape_text(sso, "SSO — un jeton JWT partagé par les trois services", font=F_BODY, size=15, color=WHITE)
# 3 apps
apps = [("ATS", "Recrutement"), ("RH · GPEC", "Mobilité interne"), ("LMS", "Formation")]
cw, gap_c, x0, top = 3.0, 0.55, 1.6, 3.15
centers = []
for i, (t, d) in enumerate(apps):
    x = x0 + i * (cw + gap_c)
    c = rect(s, x, top, cw, 1.35, fill=BLUE_LT, line_color=BLUE, line_w=1.25, rounded=True, radius=0.07)
    shape_text(c, [[ (t, {"font": F_TITLE, "size": 19, "color": INK}) ],
                   [ (d, {"font": F_BODY, "size": 13, "color": GRAY}) ]], space_after=2)
    centers.append(x + cw / 2)
    arrow(s, x + cw / 2, 2.65, x + cw / 2, top, color=BLUE, width=1.5)  # SSO -> app
# 2 bases de données
dbs = [("PostgreSQL", "RH · ATS"), ("MongoDB", "Contenus LMS")]
dy = 5.15
db_x = [centers[0] - 1.5, centers[2] - 1.5]
for (t, d), cx in zip(dbs, [centers[0], centers[2]]):
    cyl = s.shapes.add_shape(MSO_SHAPE.CAN, Inches(cx - 1.1), Inches(dy), Inches(2.2), Inches(1.1))
    cyl.shadow.inherit = False
    cyl.fill.solid(); cyl.fill.fore_color.rgb = PANEL
    cyl.line.color.rgb = GRAY_SOFT; cyl.line.width = Pt(1.0)
    shape_text(cyl, [[ (t, {"font": F_BODY, "size": 13, "color": INK, "bold": True}) ],
                     [ (d, {"font": F_BODY, "size": 11, "color": GRAY}) ]], space_after=1)
    arrow(s, cx, top + 1.35, cx, dy, color=GRAY_SOFT, width=1.25)
footer(s, 6, "Architecture")
notes(s, """
Voici la vue d'ensemble. RecruitPRO n'est pas un logiciel monolithique : ce sont trois applications spécialisées.
L'ATS pour le recrutement, le module RH et GPEC pour la mobilité interne, et le LMS pour la formation.
Ce qui les unit, c'est l'élément du haut : une authentification unique, par un jeton JWT partagé. L'utilisateur se connecte une fois et circule entre les trois sans nouvelle authentification.
Côté données : PostgreSQL pour les données structurées RH et ATS, MongoDB pour les contenus pédagogiques du LMS, qui sont plus souples.
L'idée clé à retenir : une seule identité, trois services, et donc une cohérence possible entre eux.
""")

# ===========================================================================
# SLIDE 7 — Référentiel vivant (concept)
# ===========================================================================
s = slide()
header(s, "Le référentiel vivant", "06 · Apport principal")
# noeud central
cx, cy = SW / 2, 4.35
core = rect(s, cx - 2.0, cy - 0.7, 4.0, 1.4, fill=INK, rounded=True, radius=0.12)
shape_text(core, [[ ("Référentiel", {"font": F_TITLE, "size": 19, "color": WHITE}) ],
                  [ ("de compétences vivant", {"font": F_BODY, "size": 15, "color": BLUE_LT}) ]], space_after=2)
# 3 sources
srcs = [("Recrutement", 2.4), ("Mobilité", SW / 2 - 1.4), ("Formation", SW - 5.2)]
sy = 2.15
for label, x in srcs:
    b = rect(s, x, sy, 2.8, 0.95, fill=BLUE_LT, line_color=BLUE, line_w=1.25, rounded=True, radius=0.1)
    shape_text(b, label, font=F_TITLE, size=16, color=INK)
    arrow(s, x + 1.4, sy + 0.95, cx + (x + 1.4 - cx) * 0.18, cy - 0.7, color=BLUE, width=1.75)
textbox(s, 0.9, 6.05, 11.5, 0.6, [[
    ("Trois événements alimentent ", {"color": GRAY}),
    ("une seule vérité.", {"color": BLUE, "bold": True}),
]], font=F_BODY, size=19, align=PP_ALIGN.CENTER)
footer(s, 7, "Référentiel vivant")
notes(s, """
Nous arrivons au cœur du mémoire : le référentiel vivant.
L'idée est de renverser la logique. Au lieu de mettre à jour les compétences à la main, on laisse trois événements de la vie de l'organisation les alimenter automatiquement.
Premier événement : le recrutement — quand un candidat est intégré, son profil de compétences devient celui d'un employé.
Deuxième : la mobilité interne — un changement de poste enrichit le profil.
Troisième : la formation — une compétence validée dans le LMS remonte dans le référentiel.
Ces trois flux convergent vers une source unique de vérité. Le référentiel n'est plus une photo figée : il devient un organisme qui se met à jour en permanence.
""")

# ===========================================================================
# SLIDE 8 — La boucle se referme (LMS -> RH)
# ===========================================================================
s = slide()
header(s, "La boucle se referme", "06 · Preuve du concept")
steps = [
    ("Cours validé", "dans le LMS"),
    ("Événement", "compétence émis"),
    ("Niveau + 1", "jamais abaissé"),
    ("Profil RH", "mis à jour"),
]
bw, top = 2.55, 3.1
gap_c = (SW - 1.8 - 4 * bw) / 3
x0 = 0.9
for i, (t, d) in enumerate(steps):
    x = x0 + i * (bw + gap_c)
    fill = INK if i in (0, 3) else BLUE_LT
    tcol = WHITE if i in (0, 3) else INK
    dcol = BLUE_LT if i in (0, 3) else GRAY
    b = rect(s, x, top, bw, 1.5, fill=fill, line_color=None if i in (0, 3) else BLUE,
             line_w=1.25, rounded=True, radius=0.08)
    shape_text(b, [[ (t, {"font": F_TITLE, "size": 17, "color": tcol}) ],
                   [ (d, {"font": F_BODY, "size": 12.5, "color": dcol}) ]], space_after=3)
    if i < 3:
        ax = x + bw
        arrow(s, ax, top + 0.75, ax + gap_c, top + 0.75, color=BLUE, width=2.0)
textbox(s, 0.9, 5.35, 11.5, 0.6,
        "Une compétence validée met à jour le profil — sans ressaisie.",
        font=F_BODY, size=18, color=INK, align=PP_ALIGN.CENTER, italic=True)
textbox(s, 0.9, 6.0, 11.5, 0.5,
        "Règle : le niveau ne fait que monter (course-completed).",
        font=F_BODY, size=13.5, color=GRAY, align=PP_ALIGN.CENTER)
footer(s, 8, "Boucle vivante")
notes(s, """
Prenons l'exemple le plus parlant : la formation, qui referme la boucle.
Un employé suit un cours dans le LMS et réussit son examen. À ce moment, le LMS émet un événement « compétence acquise » vers le backend RH.
Le système relève alors le niveau de la compétence correspondante dans le profil — avec une règle importante : le niveau ne fait que monter, jamais descendre. Une formation ne peut pas pénaliser quelqu'un.
Et instantanément, le profil RH est à jour. Personne n'a ressaisi quoi que ce soit.
C'est la démonstration concrète que le référentiel vit : un acte pédagogique se transforme automatiquement en donnée RH exploitable pour la mobilité ou le recrutement.
""")

# ===========================================================================
# SLIDE 9 — Moteur de matching (quadrant)
# ===========================================================================
s = slide()
header(s, "Moteur de matching", "07 · L'intelligence")
# axes
ox, oy = 4.7, 6.2          # origine
ax_w, ax_h = 6.6, 3.7
arrow(s, ox, oy, ox + ax_w, oy, color=INK, width=1.5)          # X : Fit
arrow(s, ox, oy, ox, oy - ax_h, color=INK, width=1.5)          # Y : Potential
textbox(s, ox + ax_w - 2.6, oy + 0.12, 2.7, 0.4, "Fit Score  →  besoin actuel",
        font=F_BODY, size=12.5, color=GRAY, align=PP_ALIGN.RIGHT)
tb = textbox(s, ox - 3.7, oy - ax_h - 0.15, 3.5, 0.4, "Potential Score  →  capacité d'évolution",
             font=F_BODY, size=12.5, color=GRAY, align=PP_ALIGN.RIGHT)
# pastilles recommandations
def dot(x, y, label, col):
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - 0.13), Inches(y - 0.13), Inches(0.26), Inches(0.26))
    d.shadow.inherit = False; d.fill.solid(); d.fill.fore_color.rgb = col; d.line.fill.background()
    textbox(s, x + 0.22, y - 0.22, 3.0, 0.5, label, font=F_BODY, size=13.5, color=INK, bold=True)
dot(ox + 4.9, oy - 3.1, "STRONG_FIT", BLUE)
dot(ox + 1.0, oy - 2.7, "POTENTIAL", RGBColor(0x6B, 0x46, 0xC1))
dot(ox + 0.9, oy - 0.8, "WEAK_FIT", GRAY_SOFT)
# encadré ROME
box = rect(s, 0.95, 2.05, 3.3, 2.6, fill=PANEL, line_color=LINE, line_w=1.0, rounded=True, radius=0.06)
shape_text(box, [
    [("Fit", {"font": F_TITLE, "size": 15, "color": INK})],
    [("compétences · expérience · diplôme", {"font": F_BODY, "size": 11.5, "color": GRAY})],
    [(" ", {"size": 6})],
    [("Potentiel", {"font": F_TITLE, "size": 15, "color": INK})],
    [("adjacence ROME des métiers", {"font": F_BODY, "size": 11.5, "color": GRAY})],
    [("+ bonus certification", {"font": F_BODY, "size": 11.5, "color": GRAY})],
], align=PP_ALIGN.LEFT, space_after=3)
textbox(s, 4.7, 6.55, 8.0, 0.4, "Recruter le présent, anticiper l'avenir.",
        font=F_BODY, size=14, color=BLUE, bold=True)
footer(s, 9, "Matching")
notes(s, """
Comment le système décide-t-il qu'un candidat correspond ? Avec deux scores indépendants, et c'est la deuxième contribution du travail.
Le premier, le Fit Score, mesure l'adéquation au besoin actuel : compétences, expérience, diplôme. Les compétences obligatoires pèsent davantage, et chaque compétence obligatoire manquante applique une pénalité.
Le second, le Potential Score, est plus original : il mesure la capacité à évoluer. Il s'appuie sur l'adjacence des métiers dans la nomenclature ROME — deux métiers proches partagent des compétences — et ajoute un bonus pour les certifications.
On obtient ainsi une lecture en deux dimensions : un candidat peut être « prêt maintenant » (STRONG_FIT), ou « à fort potentiel » même s'il n'est pas encore parfaitement aligné (POTENTIAL).
Concrètement, cela permet de recruter pour le présent tout en anticipant l'avenir — exactement la logique d'une GPEC.
""")

# ===========================================================================
# SLIDE 10 — Réalisation
# ===========================================================================
s = slide()
y = header(s, "Réalisation", "08 · Mise en œuvre")
bullets(s, [
    ("Backend — FastAPI · PostgreSQL", "API REST, modèle de données RH/ATS."),
    ("RH — Next.js", "Tableaux de bord, kanban, GPEC."),
    ("LMS — Next.js · MongoDB", "Cours, quiz, suivi de progression."),
    ("Intégration — SSO JWT + synchronisation", "Les trois services communiquent."),
], top=y + 0.1, size=18.5, gap=1.0)
footer(s, 10, "Réalisation")
notes(s, """
Quelques mots sur la réalisation, pour montrer que ce n'est pas une maquette mais un système fonctionnel.
Le backend est développé en FastAPI sur PostgreSQL : il porte le modèle de données et toute la logique métier — matching, promotion, GPEC.
Les deux interfaces sont en Next.js : l'application RH, avec ses tableaux de bord, son pipeline de recrutement en kanban et sa cartographie GPEC ; et le LMS, avec les cours, les quiz et le suivi de progression, adossé à MongoDB.
Enfin, l'intégration : le SSO par JWT et la synchronisation des compétences font réellement communiquer les trois services. Tout fonctionne de bout en bout, ce que la démonstration va illustrer.
""")

# ===========================================================================
# SLIDE 11 — Démonstration
# ===========================================================================
s = slide()
y = header(s, "Démonstration", "09 · En conditions réelles")
steps = [
    ("1", "Matching d'une offre", "Fit & Potentiel, analyse des écarts."),
    ("2", "Mobilité interne", "Employés proches du besoin."),
    ("3", "Assignation d'une formation", "Combler un écart ciblé."),
    ("4", "Profil RH mis à jour", "La boucle se referme."),
]
top = y + 0.05
for i, (n, t, d) in enumerate(steps):
    yy = top + i * 1.0
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), Inches(yy), Inches(0.55), Inches(0.55))
    c.shadow.inherit = False; c.fill.solid(); c.fill.fore_color.rgb = INK; c.line.fill.background()
    shape_text(c, n, font=F_TITLE, size=18, color=WHITE)
    textbox(s, 1.75, yy - 0.06, 9.8, 0.45, t, font=F_BODY, size=19, color=INK, bold=True)
    textbox(s, 1.75, yy + 0.34, 9.8, 0.4, d, font=F_BODY, size=13.5, color=GRAY)
textbox(s, 0.95, top + 4.05, 11.4, 0.4,
        "Astuce : insérer ici des captures réelles, ou faire la démonstration en direct.",
        font=F_BODY, size=12, color=GRAY_SOFT, italic=True)
footer(s, 11, "Démonstration")
notes(s, """
Je vous propose maintenant une démonstration qui suit un scénario unique, du recrutement à la mise à jour du profil.
Première étape : j'ouvre une offre et le moteur affiche les candidats classés selon les deux axes, Fit et Potentiel, avec l'analyse des écarts de compétences.
Deuxième étape : pour ce même besoin, le système propose aussi des employés en mobilité interne — des personnes déjà dans l'organisation, proches du poste.
Troisième étape : pour combler un écart précis, j'assigne une formation ciblée depuis le LMS.
Quatrième étape : une fois la formation validée, je reviens sur le profil RH — et le niveau de compétence a été relevé automatiquement. La boucle est complète.
[Si la démo est en direct : basculer vers le navigateur ici. Sinon, dérouler les captures.]
""")

# ===========================================================================
# SLIDE 12 — Résultats
# ===========================================================================
s = slide()
y = header(s, "Résultats", "10 · Ce qui est démontré")
bullets(s, [
    "Référentiel mis à jour sans ressaisie.",
    "Double score Fit / Potentiel opérationnel.",
    "Portabilité candidat → employé effective.",
    "Cartographie GPEC et analyse des écarts.",
], top=y + 0.15, size=20, gap=0.95)
footer(s, 12, "Résultats")
notes(s, """
Quels résultats ce travail démontre-t-il concrètement ?
D'abord, et c'est l'essentiel, le référentiel se met à jour sans aucune ressaisie : l'objectif central est atteint.
Ensuite, le double score Fit / Potentiel est pleinement opérationnel et change la façon de lire une candidature.
La portabilité fonctionne : un candidat recruté devient employé en conservant son profil et ses compétences.
Enfin, à l'échelle de l'organisation, le système produit une cartographie GPEC et une analyse des écarts de compétences, qui sont des outils directs de pilotage RH.
""")

# ===========================================================================
# SLIDE 13 — Limites
# ===========================================================================
s = slide()
y = header(s, "Limites", "11 · Regard critique")
bullets(s, [
    "Validation sur données simulées.",
    "L'adjacence ROME reste une approximation.",
    "Pas encore d'intégration SIRH réelle.",
    "L'évaluation pédagogique mérite d'être approfondie.",
], top=y + 0.15, size=20, gap=0.95)
footer(s, 13, "Limites")
notes(s, """
Un travail scientifique doit aussi reconnaître ses limites.
La validation a été faite sur un jeu de données simulé et cohérent, mais pas encore sur des données de production de l'ONCF.
Le Potential Score repose sur l'adjacence ROME, qui est une approximation utile mais imparfaite de la proximité réelle entre métiers.
Le système n'est pas encore connecté à un SIRH réel — paie, gestion administrative.
Enfin, la dimension d'ingénierie pédagogique — la qualité de l'évaluation des acquis dans le LMS — mérite d'être approfondie. Ce sont des limites assumées, qui ouvrent directement sur les perspectives.
""")

# ===========================================================================
# SLIDE 14 — Perspectives
# ===========================================================================
s = slide()
y = header(s, "Perspectives", "12 · Suites possibles")
bullets(s, [
    "GPEC prévisionnelle (projection des besoins).",
    "Recommandation de parcours par IA.",
    "Connecteurs vers le SIRH existant.",
    "Déploiement pilote à l'ONCF.",
], top=y + 0.15, size=20, gap=0.95)
footer(s, 14, "Perspectives")
notes(s, """
Les perspectives prolongent naturellement ces limites.
D'abord, passer d'une GPEC descriptive à une GPEC prévisionnelle, qui projette les besoins futurs en compétences — une première brique existe déjà.
Ensuite, utiliser l'intelligence artificielle pour recommander des parcours de formation individualisés à partir des écarts détectés.
Puis, développer des connecteurs vers le SIRH existant de l'ONCF pour sortir du périmètre simulé.
Et enfin, l'aboutissement logique : un déploiement pilote sur un périmètre réel de l'ONCF.
""")

# ===========================================================================
# SLIDE 15 — Conclusion + Questions
# ===========================================================================
s = slide()
rect(s, 0, 0, 0.22, SH, fill=BLUE)
textbox(s, 1.1, 2.35, 11.0, 1.6,
        "Faire vivre le référentiel\nde compétences.", font=F_LIGHT, size=40, color=INK,
        line_spacing=1.05)
hline(s, 1.13, 4.5, 4.6, color=LINE, width=1.0)
textbox(s, 1.12, 4.7, 11.0, 0.6,
        "De la donnée figée à la compétence qui évolue.", font=F_TITLE, size=18, color=BLUE)
textbox(s, 1.12, 6.2, 11.0, 0.6, [[
    ("Merci de votre attention.", {"color": INK, "bold": True}),
    ("     Vos questions.", {"color": GRAY}),
]], font=F_BODY, size=18)
notes(s, """
Pour conclure : tout ce travail tient en une phrase — faire vivre le référentiel de compétences.
RecruitPRO montre qu'en reliant recrutement, mobilité et formation autour d'une source unique de vérité, on passe d'une donnée figée à une compétence qui évolue avec la personne et avec l'organisation.
C'est, je crois, la contribution principale de ce mémoire.
Je vous remercie de votre attention et je suis maintenant à votre disposition pour vos questions.
""")

import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "RecruitPRO_Soutenance.pptx")
prs.save(out)
print("OK ->", out, "|", len(prs.slides._sldIdLst), "slides")
